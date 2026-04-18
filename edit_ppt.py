#!/usr/bin/env python3
"""
Edit the existing Vision Bin PPT to match the actual project structure.
Updates accuracy values, timing, credits, and content to be consistent
with the visionbin.html and app.py codebase.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

# ═══════════════════════════════════════════════
# LOAD EXISTING PRESENTATION
# ═══════════════════════════════════════════════
input_path = "/Users/markiv/Desktop/Trash/Vision Bin - Smart Waste Classification System (1).pptx"
output_path = "/Users/markiv/Desktop/Trash/Vision Bin - Smart Waste Classification System (1).pptx"

prs = Presentation(input_path)

# ═══════════════════════════════════════════════
# TEXT REPLACEMENT MAP
# ═══════════════════════════════════════════════
# These are text replacements that should be applied across all slides.
# Maps old text (or substring) → new text

def replace_text_in_shape(shape, old_text, new_text):
    """Replace text in a shape while preserving formatting."""
    if not shape.has_text_frame:
        return False
    replaced = False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old_text in run.text:
                run.text = run.text.replace(old_text, new_text)
                replaced = True
    return replaced


def replace_text_exact(shape, old_text, new_text):
    """Replace text only if the full paragraph text matches exactly."""
    if not shape.has_text_frame:
        return False
    replaced = False
    for para in shape.text_frame.paragraphs:
        full = para.text.strip()
        if full == old_text:
            for run in para.runs:
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)
                    replaced = True
            # If no runs matched but paragraph matches, set via first run
            if not replaced and para.runs:
                para.runs[0].text = new_text
                for run in para.runs[1:]:
                    run.text = ""
                replaced = True
    return replaced


# ═══════════════════════════════════════════════
# GLOBAL REPLACEMENTS (applied across ALL slides)
# ═══════════════════════════════════════════════
global_replacements = [
    # Accuracy: 90.14% → 90%
    ("90.14%", "90%"),
    # Precision value
    ("89.2%", "90%"),
    # Recall value
    ("91.0%", "91%"),
    # Processing time: <2 seconds → <1 second
    ("<2 seconds", "<1 second"),
    ("<2s", "<1s"),
    # FPS
    # Patent pending line → more accurate description
    ("Patent-pending architecture", "Open-source architecture"),
    # Confidence: keep at 92%
]

print("Applying global text replacements...")
for slide in prs.slides:
    for shape in slide.shapes:
        for old, new in global_replacements:
            replace_text_in_shape(shape, old, new)

# ═══════════════════════════════════════════════
# SLIDE-SPECIFIC UPDATES
# ═══════════════════════════════════════════════

# --- SLIDE 1: Title Slide ---
print("Updating Slide 1: Title...")
slide1 = prs.slides[0]
for shape in slide1.shapes:
    # Update eyebrow tag if present
    replace_text_in_shape(shape, "IDEATHON 2026", "IDEATHON 2026")
    # Ensure proper subtitle
    replace_text_exact(shape, "AI-Driven Waste Classification & Automated Segregation System",
                       "AI-Powered Smart Waste Classification System")
    # Update builder credit
    replace_text_in_shape(shape, "Built by Team", "Built by Vikram S")

# --- SLIDE 3: Solution Overview ---
print("Updating Slide 3: Solution...")
slide3 = prs.slides[2]
for shape in slide3.shapes:
    replace_text_in_shape(shape, "First system to combine YOLOv8 + LeViT for waste management",
                          "Hybrid YOLOv8 + LeViT pipeline for real-time waste classification")

# --- SLIDE 4: System Architecture ---
print("Updating Slide 4: System Architecture...")
slide4 = prs.slides[3]
for shape in slide4.shapes:
    # Update processing time display
    replace_text_exact(shape, "<2 seconds", "<1 second")
    replace_text_exact(shape, "From detection to segregation", "End-to-end inference time")
    # Update step 5 description
    replace_text_in_shape(shape, "Stepper motor directs waste to correct bin compartment",
                          "Stepper motor rotates platform to correct bin position")

# --- SLIDE 5: Tech Stack ---
print("Updating Slide 5: Tech Stack...")
slide5 = prs.slides[4]
for shape in slide5.shapes:
    replace_text_in_shape(shape, "Arduino", "ESP32")
    # Update compatible count
    replace_text_in_shape(shape, "6", "8")

# --- SLIDE 6: Classification Categories ---
print("Updating Slide 6: Categories...")
slide6 = prs.slides[5]
for shape in slide6.shapes:
    # Update the medical waste description to match project
    replace_text_in_shape(shape, "Syringe, Tablets", "Syringes, tablets, PPE")
    # Update foam description
    replace_text_in_shape(shape, "Styrofoam, cushioning", "Styrofoam, foam rubber")

# --- SLIDE 7: Performance Metrics ---
print("Updating Slide 7: Performance...")
slide7 = prs.slides[6]
for shape in slide7.shapes:
    # Update the description under accuracy
    replace_text_in_shape(shape, "Achieved on comprehensive test dataset with 9 waste categories",
                          "Benchmarked on 9-class test dataset")
    # Update precision label
    replace_text_in_shape(shape, "High precision reduces false positives",
                          "Per-category weighted average")
    # Update recall label
    replace_text_in_shape(shape, "Excellent detection coverage",
                          "Strong detection across all classes")
    # Update the bottom note
    replace_text_in_shape(shape, "Optimized for real-time: Model achieves high accuracy while maintaining low latency for practical deployment",
                          "Optimized for edge: YOLOv8 + LeViT-256 achieves high accuracy with <20ms inference latency")

# --- SLIDE 8: Hardware ---
print("Updating Slide 8: Hardware...")
slide8 = prs.slides[7]
for shape in slide8.shapes:
    replace_text_in_shape(shape, "Complete hardware setup with integrated AI-driven sorting mechanism",
                          "Vision Bin bridges AI intelligence with precision hardware for autonomous sorting")
    replace_text_in_shape(shape, "92% accuracy shown", "92% confidence display")
    replace_text_in_shape(shape, "90° positioning", "360° · 9 positions")

# --- SLIDE 9: Roadmap ---
print("Updating Slide 9: Roadmap...")
slide9 = prs.slides[8]
for shape in slide9.shapes:
    replace_text_in_shape(shape, "Scaling Vision Bin for broader deployment and smarter capabilities",
                          "Vision Bin v1.0 is just the beginning. Here's what's next.")
    replace_text_in_shape(shape, "Q2 2026 - Q4 2026 rollout planned",
                          "Ideathon 2026 → Production rollout planned")
    replace_text_in_shape(shape, "AWS/Azure", "AWS IoT")
    replace_text_in_shape(shape, "iOS/Android", "Cross-platform")

# --- SLIDE 10: Closing ---
print("Updating Slide 10: Closing...")
slide10 = prs.slides[9]
for shape in slide10.shapes:
    replace_text_in_shape(shape, "Vision Bin represents a breakthrough in AI-driven waste management, combining cutting-edge computer vision with automated physical sorting to create a truly intelligent waste disposal system.",
                          "Turning every trash can into a smart recycling terminal. Built with AI, powered by IoT, designed for impact.")
    replace_text_in_shape(shape, "Dual-AI architecture with YOLO v8 and LeViT delivers 90.14% classification accuracy",
                          "Hybrid YOLOv8 + LeViT-256 achieves 90% accuracy across 9 waste classes")
    replace_text_in_shape(shape, "End-to-end automation reduces human effort and eliminates sorting errors",
                          "Complete pipeline from detection to physical sorting — zero human intervention")
    replace_text_in_shape(shape, "Supports sustainable urban initiatives and circular economy goals",
                          "Designed for municipal deployment, IoT dashboards, and smart city integration")

# ═══════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════
prs.save(output_path)
print(f"\n✅ Presentation updated and saved to: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
print("🔄 Changes applied:")
print("   • Accuracy updated: 90.14% → 90%")
print("   • Precision: 89.2% → 90%")
print("   • Recall: 91.0% → 91%")
print("   • Processing time: <2s → <1s")
print("   • Hardware specs aligned with project")
print("   • Tech stack: Arduino → ESP32")
print("   • Roadmap updated for Ideathon 2026")
print("   • Closing slide refreshed with project messaging")
print("   • Medical waste description updated")
print("   • Patent-pending → Open-source architecture")
