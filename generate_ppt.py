#!/usr/bin/env python3
"""
VisionBin — State-of-the-Art Hackathon Presentation Generator
Generates a premium, dark-themed PowerPoint deck with custom graphics.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import math

# ═══════════════════════════════════════════════
# COLOR PALETTE
# ═══════════════════════════════════════════════
INK = RGBColor(0x0A, 0x16, 0x28)
DEEP = RGBColor(0x0D, 0x1F, 0x35)
NAVY = RGBColor(0x1B, 0x4F, 0x72)
TEAL = RGBColor(0x00, 0xA8, 0x96)
MINT = RGBColor(0x02, 0xC3, 0x9A)
CORAL = RGBColor(0xF4, 0x62, 0x3A)
GOLD = RGBColor(0xF5, 0xA6, 0x23)
WHITE = RGBColor(0xF0, 0xF6, 0xFA)
MUTED = RGBColor(0x6B, 0x8B, 0xA4)
TEXT_COLOR = RGBColor(0xC8, 0xDC, 0xE8)
BG_DARK = RGBColor(0x08, 0x11, 0x1E)
ACCENT_GREEN = RGBColor(0x4A, 0xDE, 0x80)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_gradient_bg(slide, color1, color2):
    """Set a two-stop gradient background."""
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = color2
    fill.gradient_stops[1].position = 1.0


def add_text_box(slide, text, left, top, width, height, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri', line_spacing=1.2):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox, tf


def add_multiline_text(slide, lines, left, top, width, height, default_size=16,
                       default_color=TEXT_COLOR, font_name='Calibri'):
    """Add a text box with multiple styled lines."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line_data.get('text', '')
        p.font.size = Pt(line_data.get('size', default_size))
        p.font.color.rgb = line_data.get('color', default_color)
        p.font.bold = line_data.get('bold', False)
        p.font.name = line_data.get('font', font_name)
        p.alignment = line_data.get('align', PP_ALIGN.LEFT)
        p.space_after = Pt(line_data.get('space_after', 4))
        p.space_before = Pt(line_data.get('space_before', 0))
    return txBox, tf


def add_accent_line(slide, left, top, width, height=Pt(3), color=TEAL):
    """Add a thin accent/divider line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_pill_badge(slide, text, left, top, width=Inches(2), height=Inches(0.38),
                   bg_color=TEAL, text_color=INK, font_size=10):
    """Add a rounded pill-shaped badge."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    # Adjust corner radius
    shape.adjustments[0] = 0.35
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_card(slide, left, top, width, height, bg_color=DEEP, border_color=None):
    """Add a card-like rounded rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.06
    return shape


def add_circle(slide, left, top, size, color=TEAL, opacity=1.0):
    """Add a circle shape (for decorative orbs)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Set transparency if < 1
    if opacity < 1.0:
        from lxml import etree
        spPr = shape._element.spPr
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is None:
            # Look inside ln or other places
            for sf in spPr.iter(qn('a:solidFill')):
                solidFill = sf
                break
        if solidFill is not None:
            srgb = solidFill.find(qn('a:srgbClr'))
            if srgb is not None:
                alpha = int(opacity * 100000)
                alpha_elem = srgb.find(qn('a:alpha'))
                if alpha_elem is None:
                    alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
                alpha_elem.set('val', str(alpha))
    return shape


def add_orb_decoration(slide, positions=None):
    """Add subtle decorative orbs to a slide."""
    if positions is None:
        positions = [
            (Inches(-0.8), Inches(-0.5), Inches(3), TEAL, 0.06),
            (Inches(11), Inches(5.5), Inches(2.5), CORAL, 0.04),
            (Inches(8), Inches(-1), Inches(1.8), MINT, 0.05),
        ]
    for (x, y, size, color, opacity) in positions:
        add_circle(slide, x, y, size, color, opacity)


def add_grid_decoration(slide):
    """Add subtle grid lines for tech aesthetic."""
    from lxml import etree
    for i in range(0, 14, 1):
        x = Inches(i)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(0), Pt(0.5), SLIDE_HEIGHT)
        line.fill.solid()
        line.fill.fore_color.rgb = TEAL
        line.line.fill.background()
        # Make very transparent
        spPr = line._element.spPr
        for sf in spPr.iter(qn('a:solidFill')):
            srgb = sf.find(qn('a:srgbClr'))
            if srgb is not None:
                alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
                alpha_elem.set('val', '3000')  # 3% opacity
            break


# ═══════════════════════════════════════════════
# SLIDE 1: TITLE SLIDE
# ═══════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_gradient_bg(slide1, BG_DARK, INK)
add_orb_decoration(slide1)
add_grid_decoration(slide1)

# Eyebrow
add_pill_badge(slide1, "IDEATHON 2026", Inches(0.8), Inches(1.2),
               Inches(2), Inches(0.35), TEAL, INK, 10)

# Title
add_text_box(slide1, "VISION", Inches(0.8), Inches(2.0), Inches(10), Inches(1.5),
             font_size=96, color=WHITE, bold=True, font_name='Calibri')
add_text_box(slide1, "BIN", Inches(0.8), Inches(3.3), Inches(10), Inches(1.5),
             font_size=96, color=TEAL, bold=True, font_name='Calibri')

# Accent line
add_accent_line(slide1, Inches(0.8), Inches(4.7), Inches(2), Pt(4), TEAL)

# Subtitle
add_text_box(slide1, "AI-Powered Smart Waste Classification System",
             Inches(0.8), Inches(5.0), Inches(8), Inches(0.6),
             font_size=22, color=TEXT_COLOR, bold=False)

# Description
add_text_box(slide1, "Hybrid YOLOv8 + LeViT-256  •  Real-time Edge AI  •  Autonomous Sorting",
             Inches(0.8), Inches(5.6), Inches(8), Inches(0.5),
             font_size=14, color=MUTED)

# Stats badges on the right
stats = [
    ("90%", "Accuracy"),
    ("9", "Waste Classes"),
    ("<20ms", "Inference"),
    ("100%", "Automated"),
]
for i, (val, lbl) in enumerate(stats):
    card_top = Inches(1.8) + Inches(i * 1.25)
    card = add_card(slide1, Inches(10.2), card_top, Inches(2.4), Inches(1.0), DEEP, TEAL)
    add_text_box(slide1, val, Inches(10.4), card_top + Inches(0.05), Inches(2), Inches(0.55),
                 font_size=28, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide1, lbl.upper(), Inches(10.4), card_top + Inches(0.55), Inches(2), Inches(0.35),
                 font_size=9, color=MUTED, alignment=PP_ALIGN.CENTER)

# Builder credit
add_text_box(slide1, "Built by Vikram S  •  AI / IoT Integration",
             Inches(0.8), Inches(6.6), Inches(5), Inches(0.4),
             font_size=12, color=MUTED)


# ═══════════════════════════════════════════════
# SLIDE 2: THE PROBLEM
# ═══════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(slide2, INK, DEEP)
add_orb_decoration(slide2, [
    (Inches(10), Inches(-1), Inches(3), CORAL, 0.05),
    (Inches(-1), Inches(5), Inches(2), TEAL, 0.04),
])

add_pill_badge(slide2, "// THE CHALLENGE", Inches(0.8), Inches(0.6),
               Inches(2.2), Inches(0.35), DEEP, TEAL, 10)

add_text_box(slide2, "Waste management", Inches(0.8), Inches(1.2), Inches(10), Inches(0.8),
             font_size=52, color=WHITE, bold=True)
add_text_box(slide2, "is broken.", Inches(0.8), Inches(1.95), Inches(10), Inches(0.8),
             font_size=52, color=CORAL, bold=True)

add_accent_line(slide2, Inches(0.8), Inches(2.8), Inches(1.5), Pt(3), CORAL)

# Problem cards
problems = [
    ("01", "🌍", "Environmental Damage", "Unsegregated waste contaminates soil, water, and air — accelerating climate change worldwide.", CORAL),
    ("02", "⚠️", "Inefficient Sorting", "Manual classification is slow and error-prone. Recycling rates remain critically low.", GOLD),
    ("03", "🗑", "No Intelligence", "Existing bins are passive. Without real-time classification, recyclables end up in landfill.", TEAL),
]

for i, (num, icon, title, desc, accent) in enumerate(problems):
    x = Inches(0.8) + Inches(i * 4.0)
    y = Inches(3.3)
    card = add_card(slide2, x, y, Inches(3.7), Inches(3.6), DEEP, accent)

    add_text_box(slide2, f"{num} —", x + Inches(0.3), y + Inches(0.25), Inches(3), Inches(0.3),
                 font_size=11, color=MUTED)
    add_text_box(slide2, icon, x + Inches(0.3), y + Inches(0.6), Inches(1), Inches(0.6),
                 font_size=36, color=WHITE)
    add_text_box(slide2, title, x + Inches(0.3), y + Inches(1.2), Inches(3.1), Inches(0.5),
                 font_size=20, color=WHITE, bold=True)
    add_text_box(slide2, desc, x + Inches(0.3), y + Inches(1.8), Inches(3.1), Inches(1.5),
                 font_size=13, color=MUTED, line_spacing=1.5)

# Impact stat
add_text_box(slide2, "2 billion tonnes of waste generated annually. Only 16% is recycled.",
             Inches(0.8), Inches(7.0), Inches(10), Inches(0.4),
             font_size=12, color=MUTED, alignment=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════
# SLIDE 3: OUR SOLUTION
# ═══════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(slide3, DEEP, INK)
add_orb_decoration(slide3)

add_pill_badge(slide3, "// OUR SOLUTION", Inches(0.8), Inches(0.6),
               Inches(2.2), Inches(0.35), TEAL, INK, 10)

add_text_box(slide3, "Vision Bin", Inches(0.8), Inches(1.2), Inches(10), Inches(0.9),
             font_size=56, color=WHITE, bold=True)

add_accent_line(slide3, Inches(0.8), Inches(2.2), Inches(1.5), Pt(3), TEAL)

add_text_box(slide3, "An AI-powered smart waste bin that sees, thinks, and sorts — autonomously.",
             Inches(0.8), Inches(2.5), Inches(7), Inches(0.6),
             font_size=20, color=TEXT_COLOR)

# Solution pillars
pillars = [
    ("🎯", "Real-time Detection", "YOLOv8 locates waste objects in the camera frame with bounding box precision."),
    ("🧠", "Deep Classification", "LeViT-256 Vision Transformer classifies waste into 9 categories at 90% accuracy."),
    ("⚡", "Edge Processing", "Full inference pipeline runs locally — no cloud dependency, <20ms latency."),
    ("⚙️", "Physical Sorting", "Stepper motor + rotating platform directs waste to the correct compartment automatically."),
]

for i, (icon, title, desc) in enumerate(pillars):
    x = Inches(0.8) + Inches(i * 3.1)
    y = Inches(3.5)
    card = add_card(slide3, x, y, Inches(2.85), Inches(3.3), RGBColor(0x0B, 0x19, 0x2C))

    add_text_box(slide3, icon, x + Inches(0.25), y + Inches(0.2), Inches(0.6), Inches(0.6),
                 font_size=32, color=WHITE)
    add_text_box(slide3, title, x + Inches(0.25), y + Inches(0.9), Inches(2.3), Inches(0.5),
                 font_size=16, color=WHITE, bold=True)
    add_accent_line(slide3, x + Inches(0.25), y + Inches(1.45), Inches(0.8), Pt(2), TEAL)
    add_text_box(slide3, desc, x + Inches(0.25), y + Inches(1.65), Inches(2.3), Inches(1.4),
                 font_size=12, color=MUTED, line_spacing=1.5)


# ═══════════════════════════════════════════════
# SLIDE 4: SYSTEM ARCHITECTURE / HOW IT WORKS
# ═══════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(slide4, INK, DEEP)
add_orb_decoration(slide4, [
    (Inches(6), Inches(-1), Inches(2.5), MINT, 0.04),
    (Inches(-0.5), Inches(5), Inches(2), TEAL, 0.03),
])

add_pill_badge(slide4, "// SYSTEM ARCHITECTURE", Inches(0.8), Inches(0.5),
               Inches(2.8), Inches(0.35), DEEP, TEAL, 10)

add_text_box(slide4, "Five steps to", Inches(0.8), Inches(1.1), Inches(10), Inches(0.7),
             font_size=48, color=WHITE, bold=True)
add_text_box(slide4, "a cleaner world.", Inches(0.8), Inches(1.75), Inches(10), Inches(0.7),
             font_size=48, color=TEAL, bold=True)

# Pipeline steps
steps = [
    ("01", "📷", "Capture", "Webcam feed\ncaptures waste", TEAL),
    ("02", "🎯", "Detect", "YOLOv8 locates\nobject in frame", MINT),
    ("03", "🧠", "Classify", "LeViT-256 identifies\nwaste type", TEAL),
    ("04", "🌐", "Dispatch", "Flask API sends\ncommand to IoT", MINT),
    ("05", "⚙️", "Sort", "Stepper motor\nrotates to bin", TEAL),
]

pipeline_y = Inches(3.2)
step_width = Inches(2.0)
start_x = Inches(0.8)
gap = Inches(0.4)

for i, (num, icon, title, desc, accent) in enumerate(steps):
    x = start_x + Inches(i * 2.45)

    # Circle
    circle = slide4.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.4), pipeline_y, Inches(1.2), Inches(1.2))
    circle.fill.solid()
    circle.fill.fore_color.rgb = DEEP
    circle.line.color.rgb = accent
    circle.line.width = Pt(2)

    # Icon in circle
    add_text_box(slide4, icon, x + Inches(0.4), pipeline_y + Inches(0.15), Inches(1.2), Inches(0.9),
                 font_size=36, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Step number badge
    badge = slide4.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.25), pipeline_y - Inches(0.1),
                                     Inches(0.35), Inches(0.35))
    badge.fill.solid()
    badge.fill.fore_color.rgb = accent
    badge.line.fill.background()
    add_text_box(slide4, num, x + Inches(1.22), pipeline_y - Inches(0.1), Inches(0.4), Inches(0.35),
                 font_size=9, color=INK, bold=True, alignment=PP_ALIGN.CENTER)

    # Connector arrow (except last)
    if i < 4:
        arrow_x = x + Inches(1.65)
        arrow = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, arrow_x, pipeline_y + Inches(0.55),
                                         Inches(0.8), Pt(2))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = TEAL
        arrow.line.fill.background()

        # Arrow head
        tri = slide4.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, arrow_x + Inches(0.7),
                                       pipeline_y + Inches(0.4), Inches(0.2), Inches(0.3))
        tri.fill.solid()
        tri.fill.fore_color.rgb = TEAL
        tri.line.fill.background()
        tri.rotation = 90

    # Labels
    add_text_box(slide4, title, x + Inches(0.1), pipeline_y + Inches(1.35), Inches(1.8), Inches(0.35),
                 font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide4, desc, x + Inches(0.05), pipeline_y + Inches(1.7), Inches(1.9), Inches(0.8),
                 font_size=11, color=MUTED, alignment=PP_ALIGN.CENTER, line_spacing=1.4)

# Bottom note
add_text_box(slide4, "⚡ Entire pipeline runs in under 1 second  •  No cloud dependency  •  Fully autonomous",
             Inches(0.8), Inches(6.6), Inches(11), Inches(0.4),
             font_size=13, color=TEAL, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 5: AI ENGINE - DUAL MODEL
# ═══════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(slide5, DEEP, INK)
add_orb_decoration(slide5)

add_pill_badge(slide5, "// AI ENGINE", Inches(0.8), Inches(0.5),
               Inches(1.8), Inches(0.35), DEEP, TEAL, 10)

add_text_box(slide5, "Dual-model intelligence.", Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
             font_size=48, color=WHITE, bold=True)

add_text_box(slide5, "Two state-of-the-art models work in sequence — each doing what it does best.",
             Inches(0.8), Inches(1.9), Inches(8), Inches(0.5),
             font_size=16, color=MUTED)

# YOLOv8 Card
yolo_card = add_card(slide5, Inches(0.8), Inches(2.8), Inches(5.8), Inches(4.2), RGBColor(0x0B, 0x19, 0x2C), CORAL)
add_pill_badge(slide5, "OBJECT DETECTION", Inches(1.1), Inches(3.1),
               Inches(2.2), Inches(0.3), CORAL, INK, 9)
add_text_box(slide5, "YOLOv8", Inches(1.1), Inches(3.6), Inches(5), Inches(0.7),
             font_size=40, color=WHITE, bold=True)
add_text_box(slide5, "Ultralytics  •  Real-time", Inches(1.1), Inches(4.2), Inches(5), Inches(0.3),
             font_size=12, color=MUTED)
add_text_box(slide5, "You Only Look Once — the gold standard for real-time object detection. YOLOv8 locates waste objects in the camera frame with precision bounding boxes.",
             Inches(1.1), Inches(4.6), Inches(5.2), Inches(0.8),
             font_size=13, color=TEXT_COLOR, line_spacing=1.5)

yolo_bullets = ["▸ Single-pass inference under 20ms", "▸ Anchor-free detection head",
                "▸ Optimized for edge deployment", "▸ Precision bounding box output"]
for j, bullet in enumerate(yolo_bullets):
    add_text_box(slide5, bullet, Inches(1.1), Inches(5.5) + Inches(j * 0.32), Inches(5), Inches(0.3),
                 font_size=12, color=TEAL)

# LeViT Card
levit_card = add_card(slide5, Inches(6.9), Inches(2.8), Inches(5.8), Inches(4.2), RGBColor(0x0B, 0x19, 0x2C), TEAL)
add_pill_badge(slide5, "IMAGE CLASSIFICATION", Inches(7.2), Inches(3.1),
               Inches(2.5), Inches(0.3), TEAL, INK, 9)
add_text_box(slide5, "LeViT-256", Inches(7.2), Inches(3.6), Inches(5), Inches(0.7),
             font_size=40, color=WHITE, bold=True)
add_text_box(slide5, "Vision Transformer  •  Fine-tuned", Inches(7.2), Inches(4.2), Inches(5), Inches(0.3),
             font_size=12, color=MUTED)
add_text_box(slide5, "A hybrid Vision Transformer blending attention with efficient pooling. Fine-tuned on 9 waste categories — 90% accuracy with a compact architecture.",
             Inches(7.2), Inches(4.6), Inches(5.2), Inches(0.8),
             font_size=13, color=TEXT_COLOR, line_spacing=1.5)

levit_bullets = ["▸ 90% top-1 classification accuracy", "▸ Attention + pooling hybrid layers",
                 "▸ 9-class fine-tuned waste model", "▸ Efficient transformer architecture"]
for j, bullet in enumerate(levit_bullets):
    add_text_box(slide5, bullet, Inches(7.2), Inches(5.5) + Inches(j * 0.32), Inches(5), Inches(0.3),
                 font_size=12, color=TEAL)


# ═══════════════════════════════════════════════
# SLIDE 6: PERFORMANCE / RESULTS
# ═══════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(slide6, INK, DEEP)
add_orb_decoration(slide6, [
    (Inches(0), Inches(2), Inches(3), TEAL, 0.04),
    (Inches(10), Inches(0), Inches(2), MINT, 0.03),
])

add_pill_badge(slide6, "// RESULTS", Inches(0.8), Inches(0.5),
               Inches(1.6), Inches(0.35), DEEP, TEAL, 10)

add_text_box(slide6, "Numbers that speak", Inches(0.8), Inches(1.1), Inches(10), Inches(0.7),
             font_size=48, color=WHITE, bold=True)
add_text_box(slide6, "for themselves.", Inches(0.8), Inches(1.75), Inches(10), Inches(0.7),
             font_size=48, color=TEAL, bold=True)

# Big accuracy display - left side
acc_circle = slide6.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), Inches(3.0), Inches(3.5), Inches(3.5))
acc_circle.fill.solid()
acc_circle.fill.fore_color.rgb = DEEP
acc_circle.line.color.rgb = TEAL
acc_circle.line.width = Pt(6)

add_text_box(slide6, "90%", Inches(1.2), Inches(3.8), Inches(3.5), Inches(1.0),
             font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide6, "CLASSIFICATION\nACCURACY", Inches(1.2), Inches(4.8), Inches(3.5), Inches(0.8),
             font_size=10, color=MUTED, alignment=PP_ALIGN.CENTER)

add_text_box(slide6, "BENCHMARKED ON 9-CLASS TEST DATASET",
             Inches(0.8), Inches(6.7), Inches(4), Inches(0.3),
             font_size=9, color=MUTED, alignment=PP_ALIGN.CENTER)

# Per-category bars - right side
categories = [
    ("Organic", 96), ("Cardboard", 91), ("E-Waste", 91), ("Glass", 90),
    ("Metal", 90), ("Paper", 90), ("Foam", 89), ("Plastic", 87), ("Medical", 85),
]

add_text_box(slide6, "Per-Category Precision", Inches(5.5), Inches(2.8), Inches(6), Inches(0.5),
             font_size=22, color=WHITE, bold=True)

bar_start_y = Inches(3.4)
bar_x = Inches(5.5)
label_width = Inches(1.3)
bar_width_max = Inches(4.8)
bar_height = Inches(0.28)

for i, (name, pct) in enumerate(categories):
    y = bar_start_y + Inches(i * 0.42)
    # Label
    add_text_box(slide6, name.upper(), bar_x, y, label_width, bar_height,
                 font_size=9, color=MUTED, alignment=PP_ALIGN.LEFT)
    # Track
    track = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     bar_x + label_width, y + Inches(0.05),
                                     bar_width_max, Inches(0.14), )
    track.fill.solid()
    track.fill.fore_color.rgb = RGBColor(0x15, 0x30, 0x50)
    track.line.fill.background()
    track.adjustments[0] = 0.5

    # Fill bar
    fill_width = int(bar_width_max * pct / 100)
    fill = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    bar_x + label_width, y + Inches(0.05),
                                    fill_width, Inches(0.14))
    fill.fill.solid()
    fill.fill.fore_color.rgb = TEAL
    fill.line.fill.background()
    fill.adjustments[0] = 0.5

    # Percentage
    add_text_box(slide6, f"{pct}%", bar_x + label_width + bar_width_max + Inches(0.1), y,
                 Inches(0.5), bar_height, font_size=10, color=TEAL)


# ═══════════════════════════════════════════════
# SLIDE 7: WASTE CATEGORIES
# ═══════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(slide7, DEEP, INK)
add_orb_decoration(slide7)

add_pill_badge(slide7, "// WASTE CLASSES", Inches(0.8), Inches(0.5),
               Inches(2.2), Inches(0.35), DEEP, TEAL, 10)

add_text_box(slide7, "9 categories. Zero confusion.", Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
             font_size=48, color=WHITE, bold=True)

cats = [
    ("📦", "Cardboard", "91%", GOLD),
    ("🔌", "E-Waste", "91%", CORAL),
    ("🫧", "Foam", "89%", RGBColor(0xA7, 0x8B, 0xFA)),
    ("🍶", "Glass", "90%", RGBColor(0x38, 0xBD, 0xF8)),
    ("🥫", "Metal", "90%", RGBColor(0x94, 0xA3, 0xB8)),
    ("📄", "Paper", "90%", ACCENT_GREEN),
    ("🧴", "Plastic", "87%", NAVY),
    ("🍂", "Organic", "96%", RGBColor(0x4A, 0xDE, 0x80)),
    ("💊", "Medical", "85%", RGBColor(0xF9, 0xA8, 0xD4)),
]

for i, (emoji, name, pct, dot_color) in enumerate(cats):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + Inches(col * 4.1)
    y = Inches(2.4) + Inches(row * 1.5)

    card = add_card(slide7, x, y, Inches(3.8), Inches(1.25), RGBColor(0x0B, 0x19, 0x2C), RGBColor(0x15, 0x30, 0x50))

    add_text_box(slide7, emoji, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.6),
                 font_size=28, color=WHITE)
    add_text_box(slide7, name, x + Inches(0.8), y + Inches(0.2), Inches(2), Inches(0.4),
                 font_size=16, color=WHITE, bold=True)
    add_text_box(slide7, f"{pct} precision", x + Inches(0.8), y + Inches(0.6), Inches(2), Inches(0.3),
                 font_size=11, color=MUTED)

    # Color dot
    dot = slide7.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(3.3), y + Inches(0.45),
                                   Inches(0.18), Inches(0.18))
    dot.fill.solid()
    dot.fill.fore_color.rgb = dot_color
    dot.line.fill.background()


# ═══════════════════════════════════════════════
# SLIDE 8: HARDWARE SYSTEM
# ═══════════════════════════════════════════════
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(slide8, INK, DEEP)
add_orb_decoration(slide8, [
    (Inches(10), Inches(3), Inches(3), TEAL, 0.04),
    (Inches(-1), Inches(-0.5), Inches(2.5), MINT, 0.03),
])

add_pill_badge(slide8, "// PHYSICAL SYSTEM", Inches(0.8), Inches(0.5),
               Inches(2.4), Inches(0.35), DEEP, TEAL, 10)

add_text_box(slide8, "Hardware that", Inches(0.8), Inches(1.1), Inches(10), Inches(0.7),
             font_size=48, color=WHITE, bold=True)
add_text_box(slide8, "makes it real.", Inches(0.8), Inches(1.75), Inches(10), Inches(0.7),
             font_size=48, color=TEAL, bold=True)

add_text_box(slide8, "AI alone isn't enough. Vision Bin bridges software intelligence with physical precision hardware.",
             Inches(0.8), Inches(2.5), Inches(8), Inches(0.5),
             font_size=15, color=MUTED)

# Hardware component cards
hw_components = [
    ("📷", "Camera", "HD overhead mount captures waste\nin real-time from above", "1080p · 30fps"),
    ("⚙️", "Stepper Motor", "NEMA 17 motor drives the rotating\nplatform with ±1° accuracy", "NEMA 17"),
    ("🔄", "Rotating Platform", "9-position indexed plate directs\nwaste within 500ms", "500ms · 360°"),
    ("🗑", "9 Compartments", "Dedicated, labeled bins for each\nwaste category", "Labeled · Sealed"),
]

# System diagram background
diagram_bg = add_card(slide8, Inches(0.8), Inches(3.3), Inches(11.7), Inches(3.7),
                      RGBColor(0x0B, 0x19, 0x2C), RGBColor(0x15, 0x30, 0x50))

# Header bar
header_bar = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.3),
                                      Inches(11.7), Inches(0.5))
header_bar.fill.solid()
header_bar.fill.fore_color.rgb = RGBColor(0x0E, 0x1E, 0x33)
header_bar.line.fill.background()

# Header dots
for k, dc in enumerate([CORAL, GOLD, TEAL]):
    d = slide8.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0) + Inches(k * 0.3), Inches(3.45),
                                 Inches(0.18), Inches(0.18))
    d.fill.solid()
    d.fill.fore_color.rgb = dc
    d.line.fill.background()

add_text_box(slide8, "hardware_diagram.system", Inches(2.0), Inches(3.35), Inches(3), Inches(0.4),
             font_size=10, color=MUTED)
add_text_box(slide8, "● ALL SYSTEMS OPERATIONAL", Inches(9.5), Inches(3.35), Inches(3), Inches(0.4),
             font_size=10, color=TEAL, alignment=PP_ALIGN.RIGHT)

for i, (icon, title, desc, tag) in enumerate(hw_components):
    x = Inches(1.1) + Inches(i * 2.85)
    y = Inches(4.1)

    comp = add_card(slide8, x, y, Inches(2.6), Inches(2.6), DEEP, RGBColor(0x15, 0x30, 0x50))

    add_text_box(slide8, icon, x + Inches(0.7), y + Inches(0.15), Inches(1.2), Inches(0.6),
                 font_size=32, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide8, title, x + Inches(0.15), y + Inches(0.8), Inches(2.3), Inches(0.35),
                 font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide8, desc, x + Inches(0.15), y + Inches(1.2), Inches(2.3), Inches(0.8),
                 font_size=10, color=MUTED, alignment=PP_ALIGN.CENTER, line_spacing=1.4)

    add_pill_badge(slide8, tag, x + Inches(0.5), y + Inches(2.1), Inches(1.6), Inches(0.28),
                   RGBColor(0x0B, 0x19, 0x2C), TEAL, 8)

    # Connector arrows (except last)
    if i < 3:
        arr = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(2.65), y + Inches(1.1),
                                       Inches(0.2), Pt(2))
        arr.fill.solid()
        arr.fill.fore_color.rgb = TEAL
        arr.line.fill.background()


# ═══════════════════════════════════════════════
# SLIDE 9: TECH STACK
# ═══════════════════════════════════════════════
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(slide9, DEEP, INK)
add_orb_decoration(slide9)

add_pill_badge(slide9, "// TECH STACK", Inches(0.8), Inches(0.5),
               Inches(2.0), Inches(0.35), DEEP, TEAL, 10)

add_text_box(slide9, "Built on the best tools.", Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
             font_size=48, color=WHITE, bold=True)

stack = [
    ("🐍", "Python", "Core Runtime"),
    ("🎯", "YOLOv8", "Object Detection"),
    ("🧠", "LeViT-256", "Classification"),
    ("🌐", "Flask", "REST API"),
    ("👁", "OpenCV", "Computer Vision"),
    ("⚡", "IoT Layer", "Hardware Control"),
    ("🔥", "PyTorch", "Deep Learning"),
    ("📊", "torchvision", "Transforms"),
]

for i, (icon, name, role) in enumerate(stack):
    col = i % 4
    row = i // 4
    x = Inches(0.8) + Inches(col * 3.1)
    y = Inches(2.5) + Inches(row * 2.3)

    card = add_card(slide9, x, y, Inches(2.85), Inches(2.0), RGBColor(0x0B, 0x19, 0x2C))

    add_text_box(slide9, icon, x + Inches(0.8), y + Inches(0.2), Inches(1.2), Inches(0.7),
                 font_size=36, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide9, name, x + Inches(0.2), y + Inches(1.0), Inches(2.4), Inches(0.4),
                 font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide9, role.upper(), x + Inches(0.2), y + Inches(1.4), Inches(2.4), Inches(0.3),
                 font_size=9, color=MUTED, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 10: ROADMAP / FUTURE VISION
# ═══════════════════════════════════════════════
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(slide10, INK, DEEP)
add_orb_decoration(slide10, [
    (Inches(11), Inches(0), Inches(3), TEAL, 0.04),
    (Inches(-1), Inches(4), Inches(2), MINT, 0.03),
])

add_pill_badge(slide10, "// WHAT'S NEXT", Inches(0.8), Inches(0.5),
               Inches(2.0), Inches(0.35), DEEP, TEAL, 10)

add_text_box(slide10, "The road ahead.", Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
             font_size=48, color=WHITE, bold=True)

add_text_box(slide10, "Vision Bin v1.0 is just the beginning.",
             Inches(0.8), Inches(1.9), Inches(6), Inches(0.4),
             font_size=16, color=MUTED)

roadmap = [
    ("Phase 1", "☁️", "IoT Cloud Dashboard", "Real-time bin monitoring, fill-level alerts, and waste trend analytics.", ["MQTT", "AWS IoT", "Grafana"]),
    ("Phase 2", "📱", "Mobile App", "User-facing app with eco-rewards, pickup schedules, community leaderboards.", ["React Native", "Push Alerts", "Gamification"]),
    ("Phase 3", "⚡", "Edge AI Optimization", "Full inference on Jetson Nano / Raspberry Pi — zero cloud dependency.", ["Jetson Nano", "TensorRT", "RPi"]),
    ("Phase 4", "🌆", "Smart City Analytics", "Municipal dashboards with route optimization & sustainability KPIs.", ["GIS", "Route Opt", "Open Data"]),
]

# Timeline line
timeline_x = Inches(1.5)
add_accent_line(slide10, timeline_x, Inches(2.8), Pt(2), Inches(4.2), TEAL)

for i, (phase, icon, title, desc, tags) in enumerate(roadmap):
    y = Inches(2.8) + Inches(i * 1.05)

    # Dot on timeline
    dot = slide10.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.35), y - Inches(0.05),
                                    Inches(0.3), Inches(0.3))
    dot.fill.solid()
    dot.fill.fore_color.rgb = TEAL
    dot.line.fill.background()

    add_text_box(slide10, icon, Inches(1.32), y - Inches(0.08), Inches(0.35), Inches(0.35),
                 font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Phase label
    add_text_box(slide10, phase.upper(), Inches(2.0), y - Inches(0.1), Inches(1.2), Inches(0.3),
                 font_size=10, color=TEAL, bold=True)

    # Title
    add_text_box(slide10, title, Inches(3.3), y - Inches(0.1), Inches(3.5), Inches(0.35),
                 font_size=18, color=WHITE, bold=True)

    # Description
    add_text_box(slide10, desc, Inches(6.8), y - Inches(0.1), Inches(4), Inches(0.35),
                 font_size=12, color=MUTED)

    # Tags
    for j, tag in enumerate(tags):
        add_pill_badge(slide10, tag, Inches(10.8) + Inches(j * 0) , y - Inches(0.05) + Inches(j * 0.3),
                       Inches(1.5), Inches(0.25), DEEP, MUTED, 8)


# ═══════════════════════════════════════════════
# SLIDE 11: LIVE DEMO OVERVIEW
# ═══════════════════════════════════════════════
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(slide11, DEEP, INK)
add_orb_decoration(slide11)

add_pill_badge(slide11, "// LIVE DEMO", Inches(0.8), Inches(0.5),
               Inches(1.8), Inches(0.35), TEAL, INK, 10)

add_text_box(slide11, "See the machine work.", Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
             font_size=52, color=WHITE, bold=True)

add_text_box(slide11, "Watch VisionBin detect, classify, rotate, and sort waste — fully autonomous.",
             Inches(0.8), Inches(2.0), Inches(8), Inches(0.5),
             font_size=18, color=MUTED)

# 4-step demo flow
demo_steps = [
    ("1", "🎯", "DETECT", "Camera identifies\nwaste on platform"),
    ("2", "🔄", "ROTATE", "Stepper motor\nturns to position"),
    ("3", "📤", "DROP", "Tilting mechanism\nreleases waste"),
    ("4", "✅", "SORTED", "Waste lands in\ncorrect compartment"),
]

for i, (num, icon, title, desc) in enumerate(demo_steps):
    x = Inches(0.8) + Inches(i * 3.15)
    y = Inches(3.2)

    card = add_card(slide11, x, y, Inches(2.9), Inches(3.2), RGBColor(0x0B, 0x19, 0x2C), RGBColor(0x15, 0x30, 0x50))

    # Step number at top
    num_badge = slide11.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.1), y + Inches(0.15),
                                          Inches(0.4), Inches(0.4))
    num_badge.fill.solid()
    num_badge.fill.fore_color.rgb = TEAL
    num_badge.line.fill.background()
    add_text_box(slide11, num, x + Inches(0.1), y + Inches(0.17), Inches(0.4), Inches(0.4),
                 font_size=14, color=INK, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide11, icon, x + Inches(0.8), y + Inches(0.6), Inches(1.3), Inches(0.8),
                 font_size=48, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide11, title, x + Inches(0.2), y + Inches(1.5), Inches(2.5), Inches(0.4),
                 font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_accent_line(slide11, x + Inches(0.8), y + Inches(1.95), Inches(1.3), Pt(2), TEAL)

    add_text_box(slide11, desc, x + Inches(0.2), y + Inches(2.15), Inches(2.5), Inches(0.8),
                 font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER, line_spacing=1.5)

# Bottom CTA
add_card(slide11, Inches(3.5), Inches(6.7), Inches(6.3), Inches(0.55), TEAL)
add_text_box(slide11, "▶  LAUNCH LIVE DEMO  •  localhost:5000",
             Inches(3.5), Inches(6.72), Inches(6.3), Inches(0.5),
             font_size=14, color=INK, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 12: THANK YOU / CLOSING
# ═══════════════════════════════════════════════
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(slide12, BG_DARK, INK)
add_orb_decoration(slide12, [
    (Inches(4), Inches(1), Inches(5), TEAL, 0.05),
    (Inches(8), Inches(4), Inches(3), MINT, 0.04),
    (Inches(0), Inches(3), Inches(2), CORAL, 0.03),
])
add_grid_decoration(slide12)

# Big brand
add_text_box(slide12, "VISION", Inches(0.8), Inches(1.5), Inches(12), Inches(1.5),
             font_size=96, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide12, "BIN", Inches(0.8), Inches(2.8), Inches(12), Inches(1.5),
             font_size=96, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

add_accent_line(slide12, Inches(5.5), Inches(4.3), Inches(2.3), Pt(4), TEAL)

add_text_box(slide12, "Turning every trash can into a smart recycling terminal.",
             Inches(2), Inches(4.6), Inches(9.3), Inches(0.6),
             font_size=20, color=TEXT_COLOR, alignment=PP_ALIGN.CENTER)

# Key stats row
stat_items = [("90%", "Accuracy"), ("9", "Classes"), ("<20ms", "Latency"), ("100%", "Automated")]
for i, (val, lbl) in enumerate(stat_items):
    x = Inches(2.2) + Inches(i * 2.4)
    card = add_card(slide12, x, Inches(5.4), Inches(2.1), Inches(0.9),
                    DEEP, RGBColor(0x15, 0x30, 0x50))
    add_text_box(slide12, val, x + Inches(0.1), Inches(5.42), Inches(1.9), Inches(0.5),
                 font_size=22, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide12, lbl.upper(), x + Inches(0.1), Inches(5.85), Inches(1.9), Inches(0.3),
                 font_size=8, color=MUTED, alignment=PP_ALIGN.CENTER)

# Credit
add_text_box(slide12, "Built by Vikram S  •  Ideathon 2026  •  AI / IoT Integration",
             Inches(2), Inches(6.5), Inches(9.3), Inches(0.4),
             font_size=13, color=MUTED, alignment=PP_ALIGN.CENTER)

add_text_box(slide12, "Made with ♻️ for a sustainable future",
             Inches(2), Inches(6.95), Inches(9.3), Inches(0.4),
             font_size=11, color=TEAL, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════
output_path = "/Users/markiv/Desktop/Trash/VisionBin_Hackathon_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
print("🎨 Theme: Premium dark with teal/mint accents")
